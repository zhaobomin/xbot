from __future__ import annotations

from pathlib import Path

from ..models import DocumentSection, ParsedDocument, ScannedFile


def parse_pptx(path: Path, scanned_file: ScannedFile) -> ParsedDocument:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slide_sections: list[DocumentSection] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slide_sections.append(DocumentSection(heading=f"Slide {idx}", text="\n".join(texts)))
    title = slide_sections[0].text.splitlines()[0] if slide_sections else path.stem
    return ParsedDocument(
        source_path=scanned_file.path,
        doc_type=scanned_file.doc_type,
        title=title,
        sections=slide_sections,
        modified_time=scanned_file.modified_time,
        content_hash=scanned_file.content_hash,
    )
